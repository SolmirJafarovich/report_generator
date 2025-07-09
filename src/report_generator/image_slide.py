import math
import os

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

from .utils.units import mm


def apply_text_style(paragraph, style):
    paragraph.alignment = {
        "center": PP_ALIGN.CENTER,
        "left": PP_ALIGN.LEFT,
        "right": PP_ALIGN.RIGHT,
    }.get(getattr(style, "text_align", "left"), PP_ALIGN.LEFT)

    for run in paragraph.runs:
        run.font.name = style.font_name
        run.font.size = Pt(style.font_size)
        run.font.bold = style.bold
        run.font.italic = style.italic
        run.font.color.rgb = RGBColor.from_string(style.color)


def add_image_slide(self, data):
    layout = self.prs.slide_layouts[6]  # Blank slide
    slide = self.prs.slides.add_slide(layout)

    # === Размеры слайда ===
    slide_width = self.prs.slide_width
    slide_height = self.prs.slide_height

    # === Заголовок ===
    title_height = mm(15)
    title_shape = slide.shapes.add_textbox(
        mm(10), mm(5), slide_width - mm(20), title_height
    )
    title_tf = title_shape.text_frame
    title_tf.text = data.get("title", "")
    for paragraph in title_tf.paragraphs:
        apply_text_style(paragraph, self.text_config.title)

    # === Параметры сетки ===
    images = data.get("images", [])
    num_images = len(images)

    if num_images == 0:
        return

    max_cols = 3  # максимум 3 колонки
    cols = min(num_images, max_cols)
    rows = math.ceil(num_images / cols)

    margin_x, margin_y = mm(10), title_height + mm(10)
    gap_x, gap_y = mm(10), mm(10)
    caption_height = mm(6)

    available_width = slide_width - 2 * margin_x - gap_x * (cols - 1)
    available_height = slide_height - margin_y - gap_y * (rows - 1)

    img_w = available_width / cols
    img_h = (available_height / rows) - caption_height

    # === Отрисовка изображений и подписей ===
    for idx, image in enumerate(images):
        path = image.get("path")
        caption = image.get("caption", "")
        if not os.path.exists(path):
            continue

        col = idx % cols
        row = idx // cols

        left = margin_x + col * (img_w + gap_x)
        top = margin_y + row * (img_h + caption_height + gap_y)

        # Картинка
        slide.shapes.add_picture(path, left, top, width=img_w, height=img_h)

        # Подпись
        caption_box = slide.shapes.add_textbox(left, top + img_h, img_w, caption_height)
        caption_tf = caption_box.text_frame
        caption_tf.text = caption
        for paragraph in caption_tf.paragraphs:
            apply_text_style(paragraph, self.text_config.body)
