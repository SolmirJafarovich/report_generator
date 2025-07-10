from pathlib import Path

from pptx import Presentation

from .slide_builder import SlideBuilder
from .utils.text_config import TextConfig, TextStyle


class ReportGenerator:
    def __init__(self, report_data: dict, text_config: TextConfig = None):
        self.data = report_data
        self.template_path = report_data.get("template_path")

        # Загружаем Presentation
        self.prs = self._load_template()

        # Если text_config передан явно — используем его
        self.text_config = text_config or self._load_text_config_from_data()

        # Создаём Builder с конфигом
        self.builder = SlideBuilder(self.prs, self.text_config)

    def _load_template(self):
        if self.template_path and Path(self.template_path).exists():
            return Presentation(self.template_path)
        return Presentation()

    def _load_text_config_from_data(self):
        text_cfg = self.data.get("text_config", {})
        return TextConfig(
            title=TextStyle(**text_cfg.get("title", {})),
            body=TextStyle(**text_cfg.get("body", {})),
            table_header=TextStyle(**text_cfg.get("table_header", {})),
            table_cell=TextStyle(**text_cfg.get("table_cell", {})),
            notes=TextStyle(**text_cfg.get("notes", {})),
        )

    def build(self):
        for slide_data in self.data.get("slides", []):
            self.builder.add_slide(slide_data)

    def save(self, output_path: str = "output.pptx"):
        self.prs.save(output_path)
        print(f"[✓] Презентация сохранена: {output_path}")
