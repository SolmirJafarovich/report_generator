import math
import os

from pptx.util import Pt

from .utils.placeholder_utils import get_content_placeholder_position
from .utils.table_config import (
    HEADER_ROWS,
    emu_to_pt,
    estimate_row_height,
)
from .utils.table_loader import load_table
from .utils.text_style import apply_text_style


def add_table_slide(self, data):
    layout = self.prs.slide_layouts[1]

    for table_data in data.get("tables", []):
        path = table_data.get("path")
        sheet = table_data.get("sheet")
        if not os.path.exists(path):
            print(f"[!] Таблица не найдена: {path}")
            continue

        df = load_table(path, sheet)
        rows, cols = df.shape

        # Получаем размеры плейсхолдера
        position = get_content_placeholder_position(self.prs)
        if position is None:
            print("[!] Не удалось получить размеры плейсхолдера")
            continue

        placeholder_left, placeholder_top, placeholder_width, placeholder_height = (
            position["left"],
            position["top"],
            position["width"],
            position["height"],
        )
        placeholder_height_pt = emu_to_pt(placeholder_height)

        # Получаем стиль ячеек таблицы
        cell_style = getattr(self.text_config, "table_cell", self.text_config.body)
        line_height_pt = cell_style.font_size

        # Определяем макс. число строк, которые влезают
        max_rows_per_slide = 1  # минимум
        total = line_height_pt  # заголовок
        for row in df.values:
            row_height = estimate_row_height(row, line_height_pt)
            if total + row_height > placeholder_height_pt:
                break
            total += row_height
            max_rows_per_slide += 1

        total_parts = math.ceil(rows / max_rows_per_slide)

        for part_idx in range(total_parts):
            slide = self.prs.slides.add_slide(layout)
            title_shape = slide.shapes.title
            title_text = data.get("title", "Таблица")
            if total_parts > 1:
                title_text += f" (Часть {part_idx + 1}/{total_parts})"

            title_shape.text = title_text
            for paragraph in title_shape.text_frame.paragraphs:
                apply_text_style(paragraph, self.text_config.title)

            # Подтаблица
            start = part_idx * max_rows_per_slide
            end = min(start + max_rows_per_slide, rows)
            df_chunk = df.iloc[start:end]

            # Высота таблицы в pt → emu
            table_height_pt = line_height_pt * (len(df_chunk) + HEADER_ROWS)
            table_height_emu = Pt(table_height_pt).emu
            actual_height = min(table_height_emu, placeholder_height)

            # Добавление таблицы
            table_shape = slide.shapes.add_table(
                len(df_chunk) + HEADER_ROWS,
                cols,
                placeholder_left,
                placeholder_top,
                placeholder_width,
                actual_height,
            ).table

            # Заголовки
            for col_idx, col_name in enumerate(df.columns):
                cell = table_shape.cell(0, col_idx)
                cell.text = str(col_name)
                for paragraph in cell.text_frame.paragraphs:
                    apply_text_style(paragraph, self.text_config.table_header)

            # Ячейки
            for row_idx, row in enumerate(df_chunk.values):
                for col_idx, val in enumerate(row):
                    cell = table_shape.cell(row_idx + HEADER_ROWS, col_idx)
                    cell.text = str(val)
                    for paragraph in cell.text_frame.paragraphs:
                        apply_text_style(paragraph, self.text_config.table_cell)

            # Автоширина столбцов
            col_max_char = [0] * cols
            for row in df_chunk.values:
                for col_idx, val in enumerate(row):
                    col_max_char[col_idx] = max(col_max_char[col_idx], len(str(val)))
            for col_idx, col_name in enumerate(df.columns):
                col_max_char[col_idx] = max(col_max_char[col_idx], len(str(col_name)))

            col_weights = [math.log(1 + c) for c in col_max_char]
            total_weight = sum(col_weights)
            MIN_COL_WIDTH_PCT = 0.08

            for col_idx in range(cols):
                ratio = col_weights[col_idx] / total_weight
                ratio = max(ratio, MIN_COL_WIDTH_PCT)
                table_shape.columns[col_idx].width = int(placeholder_width * ratio)
