from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt


def apply_text_style(paragraph, style):
    paragraph.alignment = {
        "center": PP_ALIGN.CENTER,
        "left": PP_ALIGN.LEFT,
        "right": PP_ALIGN.RIGHT,
    }.get(getattr(style, "text_align", "left"), PP_ALIGN.LEFT)

    for run in paragraph.runs:
        run.font.name = style.font_name
        run.font.size = Pt(style.font_size)
        run.font.bold = style.bold
        run.font.italic = style.italic
        run.font.color.rgb = RGBColor.from_string(style.color)


def add_title_slide(self, data):
    layout = self.prs.slide_layouts[0]  # Title Slide
    slide = self.prs.slides.add_slide(layout)

    # === Заголовок ===
    title_shape = slide.shapes.title
    title_shape.text = data.get("title", "")
    for paragraph in title_shape.text_frame.paragraphs:
        apply_text_style(paragraph, self.text_config.title)

    # === Подзаголовок ===
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = data.get("subtitle", "")
    for paragraph in subtitle_shape.text_frame.paragraphs:
        apply_text_style(paragraph, self.text_config.body)
