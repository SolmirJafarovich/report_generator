import copy
import logging

from pptx.enum.shapes import PP_PLACEHOLDER

logger = logging.getLogger(__name__)


def get_content_placeholder_position(prs):
    """
    Возвращает размеры плейсхолдера "Content" из layout[1], не изменяя оригинальную презентацию.
    """
    try:
        # Используем копию, чтобы не повредить оригинал
        prs_copy = copy.deepcopy(prs)
        layout = prs_copy.slide_layouts[1]
        temp_slide = prs_copy.slides.add_slide(layout)

        placeholder = None
        for shape in temp_slide.placeholders:
            if shape.placeholder_format.type == PP_PLACEHOLDER.BODY:
                placeholder = shape
                break

        if placeholder is None:
            for shape in temp_slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    placeholder = shape
                    break

        if placeholder is None:
            logger.warning("Контентный плейсхолдер не найден.")
            return None

        return {
            "left": placeholder.left,
            "top": placeholder.top,
            "width": placeholder.width,
            "height": placeholder.height,
        }

    except Exception as e:
        logger.error(f"Ошибка при получении плейсхолдера: {e}")
        return None
